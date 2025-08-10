#!/usr/bin/env python3
"""
PowerPoint Inconsistency Detector (Fixed Version)

Fixed issues:
 - Robust API response handling for Gemini 2.5 Flash
 - Better error handling for API failures
 - Improved JSON parsing with multiple fallbacks
 - Fixed NoneType errors in LLM calls
 - Enhanced retry logic and timeout handling

Usage:
    python detect_inconsistencies.py <presentation_file_or_folder> [--output json|text] [--api-key KEY] [--no-ai] [--min-confidence 0.6]

Requirements (install via pip):
    pip install python-pptx pillow google-generativeai python-dotenv pytesseract
"""

import os
import sys
import json
import re
import time
import argparse
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
from dataclasses import dataclass, asdict
from collections import defaultdict
from io import BytesIO
from difflib import SequenceMatcher
from functools import lru_cache, wraps

# Third-party imports (handle missing packages gracefully)
try:
    from pptx import Presentation
    from PIL import Image
except ImportError as e:
    print(f"Missing package: {e}. Install with: pip install python-pptx pillow")
    sys.exit(1)

# Optional libs
try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except Exception:
    TESSERACT_AVAILABLE = False

try:
    import google.generativeai as genai
    GENAI_AVAILABLE = True
except Exception:
    GENAI_AVAILABLE = False

# dotenv for convenience
try:
    from dotenv import load_dotenv
    load_dotenv()
except Exception:
    pass

# --- Data structures -------------------------------------------------------
@dataclass
class Inconsistency:
    """Represents a detected inconsistency between slides."""
    id: str
    type: str  # 'numerical', 'textual', 'temporal', 'logical', 'unknown', 'error'
    severity: str  # 'high', 'medium', 'low'
    slides: List[int]  # slide numbers involved
    description: str
    details: Dict[str, Any]
    confidence: float  # 0.0 to 1.0
    snippets: List[Dict[str, Any]] = None
    suggested_action: Optional[str] = None

# --- Utility helpers ------------------------------------------------------
def similar(a: str, b: str) -> float:
    """Calculate similarity ratio between two strings."""
    return SequenceMatcher(None, a, b).ratio()

def retry_backoff(max_attempts=3, initial_delay=1.0, backoff_factor=2.0):
    """Decorator for retry logic with exponential backoff."""
    def decorator(fn):
        @wraps(fn)
        def wrapper(*args, **kwargs):
            delay = initial_delay
            attempt = 0
            last_exception = None
            while attempt < max_attempts:
                try:
                    return fn(*args, **kwargs)
                except Exception as e:
                    last_exception = e
                    attempt += 1
                    if attempt >= max_attempts:
                        break
                    print(f"Attempt {attempt} failed: {str(e)}, retrying in {delay}s...")
                    time.sleep(delay)
                    delay *= backoff_factor
            raise last_exception
        return wrapper
    return decorator

# Numeric parsing helpers
SCALE_MAP = {'k': 1e3, 'm': 1e6, 'b': 1e9}

def parse_scaled_number(s: str) -> Optional[float]:
    """Try to parse numbers with K/M/B suffixes and commas."""
    if not s:
        return None
    s = s.strip().replace(',', '')
    # remove currency symbols for numeric parse
    s_clean = re.sub(r'[^0-9eE\.\-kKmMbB/:%]', '', s)
    # handle ratios like "3:1" or fractions "2/3" - return None
    if re.match(r'^\d+[:/]\d+', s_clean):
        return None
    m = re.match(r'([0-9]*\.?[0-9]+)\s*([kKmMbB]?)$', s_clean)
    if m:
        val = float(m.group(1))
        suf = m.group(2).lower()
        return val * SCALE_MAP.get(suf, 1)
    # try plain float
    try:
        return float(s_clean)
    except Exception:
        return None

def extract_units(token: str) -> str:
    """Extract a canonical unit token from a matched numeric expression."""
    if not token:
        return ""
    token = token.lower()
    if '%' in token:
        return '%'
    if re.search(r'₹|\$|€|£', token):
        m = re.search(r'(₹|\$|€|£)', token)
        return m.group(1)
    if re.search(r'\b(mins?|minutes?|hrs?|hours?)\b', token):
        return 'time'
    if re.search(r'[kmb]\b', token):
        m = re.search(r'([kmb])\b', token)
        return m.group(1)
    return ''

def units_match(val1: str, val2: str) -> bool:
    """Return True if two numeric tokens represent comparable units."""
    u1 = extract_units(val1)
    u2 = extract_units(val2)
    if not u1 and not u2:
        return True
    return u1 == u2

# Robust JSON extraction from LLM text
def extract_json_array_from_text(text: str) -> Optional[List[Any]]:
    """Try multiple strategies to parse a JSON array out of text returned by an LLM."""
    if not text:
        return None
    text = text.strip()
    
    # 1) direct parse
    try:
        parsed = json.loads(text)
        if isinstance(parsed, list):
            return parsed
        elif isinstance(parsed, dict):
            # wrap single object in array
            return [parsed]
    except Exception:
        pass
    
    # 2) find first [...] balanced bracket block
    bracket_count = 0
    start_idx = -1
    for i, char in enumerate(text):
        if char == '[':
            if start_idx == -1:
                start_idx = i
            bracket_count += 1
        elif char == ']':
            bracket_count -= 1
            if bracket_count == 0 and start_idx != -1:
                candidate = text[start_idx:i+1]
                try:
                    return json.loads(candidate)
                except Exception:
                    # try to repair common issues
                    repaired = candidate.replace("'", '"')
                    repaired = re.sub(r',\s*]', ']', repaired)
                    repaired = re.sub(r',\s*}', '}', repaired)
                    try:
                        return json.loads(repaired)
                    except Exception:
                        pass
                break
    
    # 3) find multiple JSON objects and wrap them
    obj_matches = re.findall(r'\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}', text)
    if obj_matches:
        try:
            arr = []
            for m in obj_matches:
                try:
                    # try direct parse
                    obj = json.loads(m)
                    arr.append(obj)
                except Exception:
                    # try repair
                    try:
                        repaired = m.replace("'", '"')
                        repaired = re.sub(r',\s*}', '}', repaired)
                        obj = json.loads(repaired)
                        arr.append(obj)
                    except Exception:
                        continue
            return arr if arr else None
        except Exception:
            pass
    
    return None

# --- Main detector class --------------------------------------------------
class InconsistencyDetector:
    def __init__(self, api_key: Optional[str] = None, min_confidence: float = 0.6, model_name: str = "gemini-2.0-flash-exp"):
        # API key handling
        self.api_key = api_key or os.getenv('GEMINI_API_KEY')
        self.min_confidence = float(min_confidence)
        self.model_name = model_name

        # Configure genai if available
        if GENAI_AVAILABLE and self.api_key:
            try:
                genai.configure(api_key=self.api_key)
                self.model = genai.GenerativeModel(self.model_name)
                print(f"Initialized Gemini model: {self.model_name}")
            except Exception as e:
                print(f"Warning: Failed to initialize Gemini model: {e}")
                self.model = None
        else:
            self.model = None
            if GENAI_AVAILABLE and not self.api_key:
                print("Warning: GEMINI_API_KEY not found. AI analysis will be skipped.")

        self.slides_content: List[Dict[str, Any]] = []
        self.inconsistencies: List[Inconsistency] = []
        self._inc_counter = 0

    def _next_id(self) -> str:
        self._inc_counter += 1
        return f"INC-{self._inc_counter:04d}"

    def extract_slide_content(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text and table content from a .pptx file."""
        slides_content = []
        try:
            prs = Presentation(file_path)
            print(f"Successfully opened PowerPoint file with {len(prs.slides)} slides")
        except Exception as e:
            print(f"Error opening PPTX file {file_path}: {e}")
            return []

        for i, slide in enumerate(prs.slides, start=1):
            slide_data = {
                'slide_number': i,
                'title': '',
                'text_content': [],
                'tables': [],
                'raw_shapes': []
            }

            for shape in slide.shapes:
                # text content
                try:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        text = shape.text.strip()
                        if shape == slide.shapes.title:
                            slide_data['title'] = text
                        else:
                            slide_data['text_content'].append(text)
                except Exception:
                    pass

                # tables
                try:
                    if getattr(shape, "has_table", False):
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        slide_data['tables'].append(table_data)
                except Exception:
                    pass

                # store raw shapes
                slide_data['raw_shapes'].append(shape)

                # OCR for embedded images
                try:
                    if hasattr(shape, 'image') and TESSERACT_AVAILABLE:
                        image = shape.image
                        blob = image.blob
                        img = Image.open(BytesIO(blob)).convert('RGB')
                        try:
                            ocr_text = pytesseract.image_to_string(img)
                            if ocr_text and ocr_text.strip():
                                slide_data['text_content'].append(f"[OCR]: {ocr_text.strip()}")
                        except Exception:
                            pass
                except Exception:
                    pass

            slides_content.append(slide_data)

        return slides_content

    def process_folder(self, folder_path: str) -> List[Dict[str, Any]]:
        """Process a folder of slide images (png/jpg)."""
        folder = Path(folder_path)
        image_files = sorted([p for p in folder.iterdir() if p.suffix.lower() in ('.png', '.jpg', '.jpeg')])
        slides = []
        
        print(f"Found {len(image_files)} image files in folder")
        
        for idx, img_file in enumerate(image_files, start=1):
            text_extracted = ""
            if TESSERACT_AVAILABLE:
                try:
                    img = Image.open(str(img_file)).convert('RGB')
                    text_extracted = pytesseract.image_to_string(img)
                except Exception as e:
                    print(f"OCR failed for {img_file}: {e}")
                    text_extracted = ""
            
            slide_data = {
                'slide_number': idx,
                'title': '',
                'text_content': [text_extracted] if text_extracted else [],
                'tables': [],
                'image_path': str(img_file),
                'source': 'image'
            }
            slides.append(slide_data)
        return slides

    # --- LLM interaction with improved error handling ----------------------
    @lru_cache(maxsize=256)
    def _cached_llm_call(self, prompt_hash: str, prompt: str) -> str:
        """Cache wrapper for LLM calls."""
        return self._llm_call(prompt)

    @retry_backoff(max_attempts=3, initial_delay=2.0, backoff_factor=2.0)
    def _llm_call(self, prompt: str) -> str:
        """Call the Gemini model and return raw text with robust error handling."""
        if not self.model:
            raise RuntimeError("Gemini model not initialized. Check API key or use --no-ai")
        
        try:
            # Generate content with timeout and safety settings
            response = self.model.generate_content(
                prompt,
                generation_config=genai.types.GenerationConfig(
                    temperature=0.1,
                    max_output_tokens=4096,
                ),
                safety_settings=[
                    {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
                    {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
                    {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
                    {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
                ]
            )
            
            # Handle different response structures
            if hasattr(response, 'text') and response.text:
                return response.text
            elif hasattr(response, 'parts') and response.parts:
                # Extract text from parts
                text_parts = []
                for part in response.parts:
                    if hasattr(part, 'text'):
                        text_parts.append(part.text)
                return ''.join(text_parts)
            elif hasattr(response, 'candidates') and response.candidates:
                # Extract from candidates
                candidate = response.candidates[0]
                if hasattr(candidate, 'content') and hasattr(candidate.content, 'parts'):
                    text_parts = []
                    for part in candidate.content.parts:
                        if hasattr(part, 'text'):
                            text_parts.append(part.text)
                    return ''.join(text_parts)
            
            # If we get here, response structure is unexpected
            raise RuntimeError(f"Unexpected response structure: {type(response)}")
            
        except Exception as e:
            print(f"LLM call failed: {str(e)}")
            raise

    def analyze_for_inconsistencies(self, slides_chunk: List[Dict[str, Any]]) -> List[Inconsistency]:
        """Use Gemini to analyze slides for inconsistencies."""
        # Build content text
        content_text = ""
        slide_numbers = []
        
        for slide in slides_chunk:
            sn = slide.get('slide_number', '?')
            slide_numbers.append(sn)
            title = slide.get('title', '')
            content = " ".join(slide.get('text_content', [])) if slide.get('text_content') else ""
            tables = slide.get('tables', [])
            
            content_text += f"\n--- SLIDE {sn} ---\n"
            if title:
                content_text += f"Title: {title}\n"
            if content:
                content_text += f"Content: {content[:1500]}...\n" if len(content) > 1500 else f"Content: {content}\n"
            if tables:
                table_str = str(tables)[:600]
                content_text += f"Tables: {table_str}...\n" if len(str(tables)) > 600 else f"Tables: {tables}\n"

        prompt = f"""You are analyzing presentation slides for inconsistencies. Find factual, numerical, textual or temporal inconsistencies across these slides.

IMPORTANT: Return ONLY a JSON array. Each object must have these exact fields:
- "type": one of "numerical", "textual", "temporal", "logical"
- "severity": one of "high", "medium", "low"  
- "slides": array of slide numbers (integers)
- "description": brief description of the inconsistency
- "confidence": number between 0.0 and 1.0
- "details": object with additional information
- "suggested_action": string with recommended action

Look for:
- Conflicting numbers, calculations, percentages
- Contradictory statements or claims
- Timeline conflicts or impossible dates
- Logic errors or missing connections

Slides content:
{content_text}

Return only the JSON array, no other text."""

        try:
            # Create a simple hash for caching
            prompt_hash = str(hash(content_text))[:16]
            raw = self._cached_llm_call(prompt_hash, prompt)
        except Exception as e:
            # LLM failed - return error inconsistency
            inc = Inconsistency(
                id=self._next_id(),
                type='error',
                severity='high',
                slides=slide_numbers,
                description=f"AI analysis failed: {str(e)}",
                details={'error_type': 'llm_failure', 'original_error': str(e)},
                confidence=0.0,
                snippets=[],
                suggested_action="Check API key, network connection, or use --no-ai mode"
            )
            return [inc]

        # Parse JSON response
        parsed = extract_json_array_from_text(raw)
        results: List[Inconsistency] = []
        
        if parsed:
            for item in parsed:
                try:
                    # Validate required fields
                    item_type = item.get('type', 'unknown')
                    severity = item.get('severity', 'medium')
                    slides = item.get('slides', slide_numbers[:2])  # fallback to first 2 slides
                    description = item.get('description', 'Inconsistency detected')
                    details = item.get('details', {})
                    confidence = float(item.get('confidence', 0.5))
                    suggested_action = item.get('suggested_action', 'Review manually')
                    
                    # Ensure slides is a list of integers
                    if not isinstance(slides, list):
                        slides = [slides] if isinstance(slides, int) else slide_numbers[:2]
                    slides = [int(s) for s in slides if isinstance(s, (int, str)) and str(s).isdigit()]
                    
                    # Clamp confidence to valid range
                    confidence = max(0.0, min(1.0, confidence))
                    
                    inc = Inconsistency(
                        id=self._next_id(),
                        type=item_type,
                        severity=severity,
                        slides=slides,
                        description=description,
                        details=details,
                        confidence=confidence,
                        snippets=[],
                        suggested_action=suggested_action
                    )
                    results.append(inc)
                except Exception as e:
                    print(f"Warning: Skipping malformed inconsistency item: {e}")
                    continue
        else:
            # No valid JSON found - create a diagnostic inconsistency
            inc = Inconsistency(
                id=self._next_id(),
                type='parsing_error',
                severity='medium',
                slides=slide_numbers,
                description="AI returned non-JSON response",
                details={'raw_response': raw[:1000] + '...' if len(raw) > 1000 else raw},
                confidence=0.3,
                snippets=[],
                suggested_action="Review raw AI response in details"
            )
            results.append(inc)

        return results

    # --- Rule-based detection ----------------------------------------------
    def detect_basic_inconsistencies(self, slides_content: List[Dict[str, Any]]) -> List[Inconsistency]:
        """Detect basic inconsistencies using rule-based heuristics."""
        inconsistencies: List[Inconsistency] = []
        numbers_by_slide = {}

        # Enhanced number patterns
        number_patterns = [
            r'(\$\s?[0-9,]+(?:\.[0-9]{1,2})?(?:[kKmMbB])?)',     # currency
            r'([0-9,]+(?:\.[0-9]{1,2})?\s*%)',                   # percentages  
            r'([0-9,]+(?:\.[0-9]{1,2})?(?:[kKmMbB])?)\s*(hours?|hrs?|mins?|minutes?)', # time
            r'\b([0-9,]+(?:\.[0-9]{1,3})?(?:[kKmMbB])?)\b',      # general numbers
        ]

        # Extract numbers from each slide
        for slide in slides_content:
            slide_num = slide.get('slide_number')
            all_text = ""
            if slide.get('title'):
                all_text += slide['title'] + " "
            if slide.get('text_content'):
                all_text += " ".join(slide['text_content']) + " "
            
            # Process tables too
            for table in slide.get('tables', []):
                for row in table:
                    all_text += " ".join(str(cell) for cell in row) + " "

            found = []
            for pat in number_patterns:
                for m in re.finditer(pat, all_text, flags=re.IGNORECASE):
                    token = m.group(1) if m.groups() else m.group(0)
                    found.append({
                        'value': token.strip(),
                        'start': m.start(),
                        'end': m.end(),
                        'context': all_text[max(0, m.start()-80):m.end()+80].strip()
                    })
            
            numbers_by_slide[slide_num] = found

        # Compare numbers across slides
        flat = []
        for s, items in numbers_by_slide.items():
            for it in items:
                flat.append((s, it))

        for idx, (s1, info1) in enumerate(flat):
            for s2, info2 in flat[idx+1:]:
                if s1 == s2:  # skip same slide
                    continue
                    
                # Skip if units don't match
                if not units_match(info1['value'], info2['value']):
                    continue
                    
                # Check context similarity
                context_sim = similar(info1['context'], info2['context'])
                
                if context_sim > 0.3 and info1['value'] != info2['value']:
                    # Parse numeric values
                    v1 = parse_scaled_number(info1['value'])
                    v2 = parse_scaled_number(info2['value'])
                    
                    confidence = 0.6
                    if v1 is not None and v2 is not None and v1 != 0 and v2 != 0:
                        rel_diff = abs(v1 - v2) / max(abs(v1), abs(v2))
                        confidence = min(0.9, 0.5 + rel_diff * 0.4)
                    
                    severity = 'high' if confidence > 0.8 else 'medium' if confidence > 0.6 else 'low'
                    
                    inc = Inconsistency(
                        id=self._next_id(),
                        type='numerical',
                        severity=severity,
                        slides=[s1, s2],
                        description=f"Numerical discrepancy: '{info1['value']}' vs '{info2['value']}'",
                        details={
                            'value1': info1['value'],
                            'value2': info2['value'],
                            'context_similarity': round(context_sim, 3),
                            'parsed_values': {'v1': v1, 'v2': v2}
                        },
                        confidence=confidence,
                        snippets=[
                            {'slide': s1, 'snippet': info1['context']},
                            {'slide': s2, 'snippet': info2['context']}
                        ],
                        suggested_action="Verify if these numbers refer to the same metric/timeframe"
                    )
                    inconsistencies.append(inc)

        # Basic textual contradictions
        contradiction_patterns = [
            ('increase', 'decrease'), ('rising', 'falling'), ('up', 'down'),
            ('high', 'low'), ('many', 'few'), ('growing', 'shrinking'),
            ('competitive', 'monopoly'), ('strong', 'weak')
        ]
        
        slide_texts = {}
        for slide in slides_content:
            sn = slide['slide_number']
            text = " ".join([slide.get('title', '')] + slide.get('text_content', [])).lower()
            slide_texts[sn] = text

        slide_nums = list(slide_texts.keys())
        for i in range(len(slide_nums)):
            for j in range(i+1, len(slide_nums)):
                s1, s2 = slide_nums[i], slide_nums[j]
                text1, text2 = slide_texts[s1], slide_texts[s2]
                
                for word1, word2 in contradiction_patterns:
                    if ((word1 in text1 and word2 in text2) or 
                        (word2 in text1 and word1 in text2)):
                        
                        inc = Inconsistency(
                            id=self._next_id(),
                            type='textual',
                            severity='low',
                            slides=[s1, s2],
                            description=f"Potential contradiction: '{word1}' vs '{word2}'",
                            details={'contradiction_pair': [word1, word2]},
                            confidence=0.4,
                            snippets=[
                                {'slide': s1, 'snippet': text1[:200]},
                                {'slide': s2, 'snippet': text2[:200]}
                            ],
                            suggested_action="Check if statements refer to different contexts"
                        )
                        inconsistencies.append(inc)
                        break

        return inconsistencies

    # --- Main analysis orchestration ---------------------------------------
    def analyze_presentation(self, input_path: str, use_ai: bool = True) -> List[Inconsistency]:
        """Main analysis function."""
        input_path = Path(input_path)
        
        if input_path.is_file() and input_path.suffix.lower() == '.pptx':
            print(f"Extracting content from PowerPoint: {input_path}")
            self.slides_content = self.extract_slide_content(str(input_path))
        elif input_path.is_dir():
            print(f"Processing folder of images: {input_path}")
            self.slides_content = self.process_folder(str(input_path))
        else:
            raise ValueError("Input must be a .pptx file or folder with slide images")

        if not self.slides_content:
            raise ValueError("No slide content extracted")

        print(f"Extracted content from {len(self.slides_content)} slides")

        all_inconsistencies: List[Inconsistency] = []

        # Rule-based analysis
        print("Running rule-based inconsistency detection...")
        rule_results = self.detect_basic_inconsistencies(self.slides_content)
        all_inconsistencies.extend(rule_results)
        print(f"Rule-based analysis found {len(rule_results)} potential issues")

        # AI analysis if available
        if use_ai and self.model:
            print("Running AI-powered analysis...")
            chunk_size = 6  # smaller chunks for reliability
            overlap = 1
            i = 0
            ai_results = []
            
            while i < len(self.slides_content):
                chunk = self.slides_content[i:i+chunk_size]
                chunk_slides = [s['slide_number'] for s in chunk]
                print(f"  Analyzing slides {chunk_slides}...")
                
                try:
                    chunk_results = self.analyze_for_inconsistencies(chunk)
                    ai_results.extend(chunk_results)
                except Exception as e:
                    print(f"  Warning: AI analysis failed for chunk {chunk_slides}: {e}")
                
                i += max(1, chunk_size - overlap)
            
            all_inconsistencies.extend(ai_results)
            print(f"AI analysis found {len(ai_results)} potential issues")
            
        elif use_ai and not self.model:
            print("Warning: AI analysis requested but model not available")

        # Deduplication and filtering
        unique_inconsistencies = []
        for inc in all_inconsistencies:
            # Apply confidence threshold (but keep error types)
            if inc.confidence < self.min_confidence and inc.type not in ['error', 'parsing_error']:
                continue
                
            # Check for duplicates
            is_duplicate = False
            for existing in unique_inconsistencies:
                if (existing.type == inc.type and 
                    set(existing.slides) == set(inc.slides) and
                    similar(existing.description, inc.description) > 0.8):
                    # Merge by keeping higher confidence
                    existing.confidence = max(existing.confidence, inc.confidence)
                    if inc.snippets:
                        existing.snippets = (existing.snippets or []) + inc.snippets
                    is_duplicate = True
                    break
            
            if not is_duplicate:
                unique_inconsistencies.append(inc)

        # Sort by confidence and severity
        severity_order = {'high': 3, 'medium': 2, 'low': 1}
        self.inconsistencies = sorted(
            unique_inconsistencies, 
            key=lambda x: (-x.confidence, -severity_order.get(x.severity, 0))
        )
        
        print(f"Total unique inconsistencies after filtering: {len(self.inconsistencies)}")
        return self.inconsistencies

    # --- Report generation -------------------------------------------------
    def generate_report(self, output_format: str = 'text') -> str:
        """Generate formatted report of inconsistencies."""
        if not self.inconsistencies:
            return "No inconsistencies detected above the confidence threshold."

        if output_format == 'json':
            data = {
                'metadata': {
                    'total_slides': len(self.slides_content),
                    'total_inconsistencies': len(self.inconsistencies),
                    'confidence_threshold': self.min_confidence,
                    'analysis_timestamp': time.strftime('%Y-%m-%d %H:%M:%S')
                },
                'inconsistencies': [asdict(inc) for inc in self.inconsistencies]
            }
            return json.dumps(data, indent=2, ensure_ascii=False)

        # Text format report
        lines = []
        lines.append("POWERPOINT INCONSISTENCY ANALYSIS REPORT")
        lines.append("=" * 60)
        lines.append(f"Analysis Date: {time.strftime('%Y-%m-%d %H:%M:%S')}")
        lines.append(f"Total Slides: {len(self.slides_content)}")
        lines.append(f"Inconsistencies Found: {len(self.inconsistencies)}")
        lines.append(f"Confidence Threshold: {self.min_confidence:.1%}")
        lines.append("")

        # Group by severity
        by_severity = defaultdict(list)
        for inc in self.inconsistencies:
            by_severity[inc.severity].append(inc)

        for severity in ['high', 'medium', 'low']:
            issues = by_severity.get(severity, [])
            if not issues:
                continue
                
            lines.append(f"{severity.upper()} SEVERITY ISSUES ({len(issues)})")
            lines.append("-" * 50)
            
            for i, inc in enumerate(issues, 1):
                lines.append(f"\n{i}. [{inc.id}] {inc.type.upper()} INCONSISTENCY")
                lines.append(f"   Slides: {', '.join(map(str, inc.slides))}")
                lines.append(f"   Confidence: {inc.confidence:.1%}")
                lines.append(f"   Description: {inc.description}")
                
                if inc.suggested_action:
                    lines.append(f"   Suggested Action: {inc.suggested_action}")
                
                if inc.details:
                    lines.append("   Details:")
                    for key, value in inc.details.items():
                        value_str = str(value)
                        if len(value_str) > 150:
                            value_str = value_str[:150] + "..."
                        lines.append(f"     {key}: {value_str}")
                
                if inc.snippets:
                    lines.append("   Evidence:")
                    for snippet in inc.snippets[:3]:  # limit to 3 snippets
                        slide_num = snippet.get('slide', '?')
                        text = snippet.get('snippet', '')[:200]
                        lines.append(f"     Slide {slide_num}: \"{text}...\"")
            
            lines.append("\n" + "=" * 60)

        return "\n".join(lines)


# --- Command Line Interface --------------------------------------------
def main():
    parser = argparse.ArgumentParser(
        description='Detect inconsistencies in PowerPoint presentations using AI and rule-based analysis',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python detect_inconsistencies.py presentation.pptx
  python detect_inconsistencies.py slides_folder/ --no-ai
  python detect_inconsistencies.py deck.pptx --output json --min-confidence 0.8
  python detect_inconsistencies.py file.pptx --api-key YOUR_GEMINI_KEY
        """
    )
    
    parser.add_argument('input_path', 
                       help='Path to .pptx file or folder containing slide images')
    parser.add_argument('--output', '-o', 
                       choices=['text', 'json'], 
                       default='text',
                       help='Output format (default: text)')
    parser.add_argument('--api-key', 
                       help='Gemini API key (or set GEMINI_API_KEY environment variable)')
    parser.add_argument('--no-ai', 
                       action='store_true',
                       help='Skip AI analysis and use only rule-based detection')
    parser.add_argument('--min-confidence', 
                       type=float, 
                       default=0.5,
                       help='Minimum confidence threshold (0.0-1.0, default: 0.5)')
    parser.add_argument('--model',
                       default='gemini-2.0-flash-exp',
                       help='Gemini model to use (default: gemini-2.0-flash-exp)')

    args = parser.parse_args()

    # Validate arguments
    if not Path(args.input_path).exists():
        print(f"Error: Input path '{args.input_path}' does not exist")
        sys.exit(1)
        
    if not 0.0 <= args.min_confidence <= 1.0:
        print("Error: min-confidence must be between 0.0 and 1.0")
        sys.exit(1)

    # Initialize detector
    try:
        detector = InconsistencyDetector(
            api_key=args.api_key,
            min_confidence=args.min_confidence,
            model_name=args.model
        )
    except Exception as e:
        print(f"Error initializing detector: {e}")
        if args.no_ai:
            print("Continuing with rule-based analysis only...")
            detector = InconsistencyDetector(api_key=None, min_confidence=args.min_confidence)
        else:
            print("Use --no-ai to skip AI analysis, or check your API key")
            sys.exit(1)

    # Run analysis
    use_ai = not args.no_ai
    if use_ai and not GENAI_AVAILABLE:
        print("Warning: google-generativeai not installed. Running rule-based analysis only.")
        use_ai = False
    elif use_ai and not detector.model:
        print("Warning: AI model not available. Running rule-based analysis only.")
        use_ai = False

    try:
        print(f"Starting analysis of: {args.input_path}")
        inconsistencies = detector.analyze_presentation(args.input_path, use_ai=use_ai)
        
        # Generate and display report
        report = detector.generate_report(output_format=args.output)
        print("\n" + "="*80)
        print(report)
        
        # Save report to file
        timestamp = time.strftime('%Y%m%d_%H%M%S')
        output_ext = 'json' if args.output == 'json' else 'txt'
        output_filename = f"inconsistency_report_{timestamp}.{output_ext}"
        
        with open(output_filename, 'w', encoding='utf-8') as f:
            f.write(report)
        
        print(f"\nReport saved to: {output_filename}")
        
        # Summary
        if inconsistencies:
            high = sum(1 for inc in inconsistencies if inc.severity == 'high')
            medium = sum(1 for inc in inconsistencies if inc.severity == 'medium')
            low = sum(1 for inc in inconsistencies if inc.severity == 'low')
            print(f"\nSummary: {high} high, {medium} medium, {low} low severity issues found")
        else:
            print("\nNo inconsistencies found above the confidence threshold.")
            
    except KeyboardInterrupt:
        print("\nAnalysis interrupted by user")
        sys.exit(1)
    except Exception as e:
        print(f"\nError during analysis: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()